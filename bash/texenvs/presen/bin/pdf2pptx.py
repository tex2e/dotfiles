#!/usr/bin/env python3

import os
import sys
import glob
from pptx import Presentation # pip install python-pptx

if len(sys.argv) < 1:
    print("Usage: pdf2pptx <path/to/pages/> [path/to/presen-note.txt]")
    sys.exit()

def listdir_nohidden(path):
    return glob.glob(os.path.join(path, '*'))

def main(argv):
    images_path = argv[1]
    note_path = argv[2] if len(argv) > 2 else None
    notes = None
    if note_path:
        with open(note_path) as f:
            notes = f.read().strip().split("\n\n")

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    left   = 0
    top    = 0
    width  = prs.slide_width
    height = prs.slide_height

    for i, image_path in enumerate(sorted(listdir_nohidden(images_path))):
        print('Create:', image_path)
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image_path, left, top, width, height)

        if len(notes) > i:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes[i]

    prs.save('presen.pptx')
    print('Craete: presen.pptx')

if __name__ == '__main__':
    main(sys.argv)
